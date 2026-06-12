import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

import argparse
import shutil
import json
from calendar import monthrange
from datetime import datetime, timedelta
from dateutil.relativedelta import relativedelta
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image
from core.generate_commentary import generate_monthly_slide_content, generate_mtd_slide_content
from core.get_funnel_data import fmt_int, fmt_pct, fmt_gbp
from monthly_reports.generate_visualisation import render_graph, initialise_brand, BRAND
from monthly_reports.generate_data_export import export_slide_data

initialise_brand()

# ── SLIDE LAYOUT INDICES ──────────────────────────────────────────────────────
SLD_LAYOUT_COVER                      = 0
SLD_LAYOUT_SECTION_SEPARATOR          = 1
SLD_LAYOUT_SUNSET_SECTION_SEPARATOR   = 2
SLD_LAYOUT_BURNT_ORANGE_SECTION_SEPARATOR = 3
SLD_LAYOUT_QUOTATION                  = 4
SLD_LAYOUT_TITLE_AND_BODY             = 5
SLD_LAYOUT_TITLE_AND_BODY_1           = 6
SLD_LAYOUT_TITLE_AND_BODY_2           = 7
SLD_LAYOUT_TITLE_AND_BODY_3           = 8
SLD_LAYOUT_2_COLUMN                   = 9
SLD_LAYOUT_3_COLUMN                   = 10
SLD_LAYOUT_4_COLUMN                   = 11
SLD_LAYOUT_BIG_NUMBER                 = 12
SLD_LAYOUT_BLANK                      = 13

# ── BRAND COLOURS ─────────────────────────────────────────────────────────────
C = {
    "gold":   RGBColor(0xFE, 0xC0, 0x42),
    "orange": RGBColor(0xF2, 0x7D, 0x39),
    "teal":   RGBColor(0x4F, 0xA6, 0xA4),
    "dark":   RGBColor(0x2B, 0x2D, 0x42),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "light":  RGBColor(0xFF, 0xF7, 0xE4),
    "grey":   RGBColor(0xD9, 0xD9, 0xD9),
}

STATUS_COLOURS = {
    "Complete":    C["teal"],
    "In Progress": C["gold"],
    "Scheduled":   C["grey"],
    "Blocked":     C["orange"],
}

HEADER_TEXT_COLOUR = {
    "Complete":    C["white"],
    "In Progress": C["dark"],
    "Scheduled":   C["dark"],
    "Blocked":     C["white"],
}


# ── PRIVATE HELPERS ───────────────────────────────────────────────────────────

def _set_text(tf, text, size=Pt(14)):
    tf.text = text
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = size


def _shrink_bullet(p, pct=70):
    pPr = p._p.get_or_add_pPr()
    for el in pPr.findall(qn('a:buSzPct')):
        pPr.remove(el)
    bu = pPr.makeelement(qn('a:buSzPct'), attrib={'val': str(pct * 1000)})
    pPr.append(bu)


def _populate_bullets(tf, bullets, size=Pt(12)):
    tf.clear()
    for i, bullet in enumerate(bullets):
        text = bullet['point'] if isinstance(bullet, dict) else bullet
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        _shrink_bullet(p)
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = size


def _add_chart_image(slide, chart_path, left=Inches(5), top=Inches(1.5), width=Inches(4.5)):
    with Image.open(chart_path) as img:
        img_w, img_h = img.size
    height = int(width * img_h / img_w)
    slide.shapes.add_picture(chart_path, left=left, top=top, width=width, height=height)


def _add_title_textbox(slide, title):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = BRAND["font"]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = C["dark"]


def _style_cell(cell, text, bg_colour=None, text_colour=None, bold=False,
                size=Pt(10), align=PP_ALIGN.LEFT):
    if bg_colour:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_colour
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for run in p.runs:
        run.font.name = BRAND["font"]
        run.font.size = size
        run.font.bold = bold
        if text_colour:
            run.font.color.rgb = text_colour


def _add_table_shape(slide, headers, rows, left, top, width, height, status_col=None, totals_row=None):
    extra = 1 if totals_row else 0
    n_rows = len(rows) + 1 + extra
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    for j, header in enumerate(headers):
        _style_cell(tbl.cell(0, j), header,
                    bg_colour=C["dark"], text_colour=C["white"],
                    bold=True, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        row_bg = C["light"] if i % 2 == 0 else C["white"]
        for j, val in enumerate(row):
            cell_bg = STATUS_COLOURS.get(val, row_bg) if (status_col is not None and j == status_col) else row_bg
            _style_cell(tbl.cell(i + 1, j), val, bg_colour=cell_bg, text_colour=C["dark"], align=PP_ALIGN.CENTER)

    if totals_row:
        totals_row_idx = len(rows) + 1
        for j, val in enumerate(totals_row):
            _style_cell(tbl.cell(totals_row_idx, j), val,
                        bg_colour=C["dark"], text_colour=C["white"],
                        bold=True, align=PP_ALIGN.CENTER)

    return tbl


def _add_kpi_boxes(slide, kpis, start_x=Inches(7.1), start_y=None,
                   box_w=Inches(2.5), box_h=Inches(0.95), gap=Inches(0.1)):
    if start_y is None:
        n = len(kpis)
        total_h = n * box_h + (n - 1) * gap
        # Body content area on TITLE_AND_BODY layout: top=1.30", bottom=4.62"
        body_top    = Inches(1.30)
        body_bottom = Inches(4.62)
        start_y = body_top + (body_bottom - body_top - total_h) / 2

    box_colours = [C["gold"], C["orange"], C["teal"], C["dark"]]

    for i, (label, data) in enumerate(kpis):
        y = start_y + i * (box_h + gap)
        shape = slide.shapes.add_shape(5, start_x, y, box_w, box_h)
        shape.adjustments[0] = 0.08
        shape.fill.solid()
        shape.fill.fore_color.rgb = box_colours[i % len(box_colours)]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_top    = Inches(0.1)
        tf.margin_bottom = Inches(0.08)
        tf.margin_left   = Inches(0.1)
        tf.margin_right  = Inches(0.1)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = BRAND["font"]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = C["dark"]

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = data.get('curr', '—')
        run2.font.name = BRAND["font"]
        run2.font.size = Pt(17)
        run2.font.bold = True
        run2.font.color.rgb = C["dark"]

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = f"vs {data.get('prev', '—')}  ({data.get('pct', '—')})"
        run3.font.name = BRAND["font"]
        run3.font.size = Pt(9)
        run3.font.color.rgb = C["dark"]


def _add_kpi_boxes_horizontal(slide, prs, kpis, top=None, box_h=Inches(1.1), gap=Inches(0.15)):
    n = len(kpis)
    margin = Inches(0.5)
    total_w = prs.slide_width - 2 * margin
    box_w = (total_w - (n - 1) * gap) // n
    if top is None:
        top = prs.slide_height - box_h - Inches(0.25)
    box_colours = [C["gold"], C["orange"], C["teal"], C["dark"]]

    for i, (label, data) in enumerate(kpis):
        x = margin + i * (box_w + gap)
        shape = slide.shapes.add_shape(5, int(x), int(top), int(box_w), int(box_h))
        shape.adjustments[0] = 0.08
        shape.fill.solid()
        shape.fill.fore_color.rgb = box_colours[i % len(box_colours)]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_top    = Inches(0.08)
        tf.margin_bottom = Inches(0.06)
        tf.margin_left   = Inches(0.1)
        tf.margin_right  = Inches(0.1)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = BRAND["font"]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = C["dark"]

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = data.get('curr', '—')
        run2.font.name = BRAND["font"]
        run2.font.size = Pt(18)
        run2.font.bold = True
        run2.font.color.rgb = C["dark"]

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = f"vs {data.get('prev', '—')}  ({data.get('pct', '—')})"
        run3.font.name = BRAND["font"]
        run3.font.size = Pt(9)
        run3.font.color.rgb = C["dark"]


def _build_kpis_for(client, data_key='paid_data', kpi_count=None):
    paid_total = client.get(data_key, {}).get('Total', {})
    if client.get('account_type') == 'Lead Gen':
        all_kpis = [
            ('Cost',            paid_total.get('Cost', {})),
            ('Conversions',     paid_total.get('Conversions', {})),
            ('CPA',             paid_total.get('CPA', {})),
            ('Conversion Rate', paid_total.get('Conversion Rate', {})),
        ]
    else:
        all_kpis = [
            ('Cost',            paid_total.get('Cost', {})),
            ('Revenue',         paid_total.get('Transaction Revenue', {})),
            ('ROAS',            paid_total.get('ROAS', {})),
            ('Conversion Rate', paid_total.get('Conversion Rate', {})),
        ]
    if kpi_count is not None:
        return all_kpis[:max(1, min(4, int(kpi_count)))]
    return all_kpis[:3]


def _extract_current_tasks(plan_json):
    if not plan_json:
        return []
    if isinstance(plan_json, list):
        return plan_json
    if 'tasks' in plan_json and 'plan_status' not in plan_json:
        return plan_json['tasks']

    tasks = []
    for quarter_data in plan_json.values():
        if isinstance(quarter_data, dict) and quarter_data.get('plan_status') == 'current':
            tasks = quarter_data.get('tasks', [])
            if not isinstance(tasks, list):
                return []
            break

    today = datetime.now()
    anchor = today.replace(day=1) - relativedelta(months=1) if today.day <= 5 else today.replace(day=1)
    month_start = anchor
    month_end = (anchor + relativedelta(months=1)) - timedelta(days=1)

    def _overlaps(t):
        for fmt in ('%d/%m/%y', '%d/%m/%Y', '%Y-%m-%d'):
            try:
                s = datetime.strptime(t.get('start_date', ''), fmt)
                e = datetime.strptime(t.get('end_date', ''), fmt)
                return s <= month_end and e >= month_start
            except (ValueError, TypeError):
                continue
        return False

    return [t for t in tasks
            if t.get('category', '').strip().lower() == 'active workstream'
            and _overlaps(t)]


def _extract_all_plan_tasks(plan_json):
    """Returns (tasks, plan_start_str, plan_end_str) for the current quarter."""
    if not plan_json:
        return [], None, None

    def _active(tasks):
        return [t for t in tasks if t.get('category', '').strip().lower() == 'active workstream']

    if isinstance(plan_json, list):
        return _active(plan_json), None, None
    if 'tasks' in plan_json and 'plan_status' not in plan_json:
        return _active(plan_json['tasks']), None, None
    for quarter_data in plan_json.values():
        if isinstance(quarter_data, dict) and quarter_data.get('plan_status') == 'current':
            tasks = quarter_data.get('tasks', [])
            return (
                _active(tasks) if isinstance(tasks, list) else [],
                quarter_data.get('plan_start'),
                quarter_data.get('plan_end'),
            )
    return [], None, None


def _fmt_date(date_str):
    if not date_str:
        return ''
    for fmt in ('%d/%m/%y', '%d/%m/%Y', '%Y-%m-%d'):
        try:
            return datetime.strptime(date_str, fmt).strftime('%d/%m/%Y')
        except ValueError:
            continue
    return date_str


def _parse_date_for_sort(date_str):
    if not date_str:
        return datetime.max
    for fmt in ('%d/%m/%y', '%d/%m/%Y', '%Y-%m-%d'):
        try:
            return datetime.strptime(date_str, fmt)
        except ValueError:
            continue
    return datetime.max


# ── PUBLIC SLIDE TEMPLATES ────────────────────────────────────────────────────

def slide_cover(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_COVER])
    slide.placeholders[0].text = title
    return slide


def slide_section_separator(prs, title, variant='navy'):
    layout_map = {
        'navy':   SLD_LAYOUT_SECTION_SEPARATOR,
        'gold':   SLD_LAYOUT_SUNSET_SECTION_SEPARATOR,
        'orange': SLD_LAYOUT_BURNT_ORANGE_SECTION_SEPARATOR,
    }
    slide = prs.slides.add_slide(prs.slide_layouts[layout_map.get(variant, SLD_LAYOUT_SECTION_SEPARATOR)])
    slide.placeholders[0].text = title
    return slide


def slide_commentary(prs, title, summary, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_BODY])
    slide.placeholders[0].text = title
    try:
        _set_text(slide.placeholders[4].text_frame, summary)
    except (KeyError, IndexError):
        pass
    try:
        _populate_bullets(slide.placeholders[1].text_frame, bullets)
    except (KeyError, IndexError):
        pass
    return slide


def _iso_to_dmy(iso_str):
    try:
        return datetime.strptime(iso_str[:10], "%Y-%m-%d").strftime("%d/%m/%Y")
    except Exception:
        return iso_str


def _fmt_date_label(start, end, comp_start=None, comp_end=None):
    label = f"{start} - {end}"
    if comp_start and comp_end:
        label += f" vs {comp_start} - {comp_end}"
    return label + " | Ad Platform & GA4"


def slide_chart_commentary(prs, title, summary, bullets, chart_path, date_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_BODY])
    slide.placeholders[0].text = title
    try:
        if date_label:
            slide.placeholders[2].text = date_label
    except (KeyError, IndexError):
        pass
    try:
        _set_text(slide.placeholders[4].text_frame, summary)
    except (KeyError, IndexError):
        pass
    try:
        _populate_bullets(slide.placeholders[1].text_frame, bullets)
    except (KeyError, IndexError):
        pass
    _add_chart_image(slide, chart_path)
    return slide


def slide_scorecard_commentary(prs, title, summary, bullets, kpis, date_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_BODY])
    slide.placeholders[0].text = title
    try:
        if date_label:
            slide.placeholders[2].text = date_label
    except (KeyError, IndexError):
        pass
    try:
        _set_text(slide.placeholders[4].text_frame, summary)
    except (KeyError, IndexError):
        pass
    try:
        _populate_bullets(slide.placeholders[1].text_frame, bullets)
    except (KeyError, IndexError):
        pass
    _add_kpi_boxes(slide, kpis)
    return slide


def slide_scorecard_vertical(prs, title, summary, bullets, kpis, date_label=None):
    return slide_scorecard_commentary(prs, title, summary, bullets, kpis, date_label)


def slide_scorecard_horizontal(prs, title, summary, bullets, kpis, date_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_BODY])
    slide.placeholders[0].text = title
    try:
        if date_label:
            slide.placeholders[2].text = date_label
    except (KeyError, IndexError):
        pass
    try:
        _set_text(slide.placeholders[4].text_frame, summary)
    except (KeyError, IndexError):
        pass
    # Constrain bullets to upper portion — KPI row occupies bottom ~1.4"
    try:
        ph = slide.placeholders[1]
        ph.top    = Inches(1.30)
        ph.height = Inches(3.20)
        _populate_bullets(ph.text_frame, bullets)
    except (KeyError, IndexError):
        pass
    _add_kpi_boxes_horizontal(slide, prs, kpis)
    return slide


def slide_full_chart(prs, title, summary, chart_path, date_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_BODY])
    slide.placeholders[0].text = title
    try:
        if date_label:
            slide.placeholders[2].text = date_label
    except (KeyError, IndexError):
        pass
    try:
        _set_text(slide.placeholders[4].text_frame, summary)
    except (KeyError, IndexError):
        pass
    with Image.open(chart_path) as img:
        img_w, img_h = img.size
    chart_top    = Inches(1.55)
    chart_width  = prs.slide_width - Inches(0.8)
    chart_height = int(chart_width * img_h / img_w)
    max_h        = prs.slide_height - chart_top - Inches(0.2)
    if chart_height > max_h:
        chart_height = max_h
        chart_width  = int(chart_height * img_w / img_h)
    left = (prs.slide_width - chart_width) // 2
    slide.shapes.add_picture(chart_path, left=left, top=chart_top,
                             width=chart_width, height=chart_height)
    return slide


def slide_big_number(prs, title, summary, chart_path,
                     hero_label, hero_curr, hero_prev, hero_pct, date_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_BODY])
    slide.placeholders[0].text = title
    try:
        if date_label:
            slide.placeholders[2].text = date_label
    except (KeyError, IndexError):
        pass
    try:
        _set_text(slide.placeholders[4].text_frame, summary)
    except (KeyError, IndexError):
        pass

    box_top  = Inches(1.55)
    box_h    = Inches(3.1)
    box_w    = Inches(4.0)

    shape = slide.shapes.add_shape(5, Inches(0.4), box_top, box_w, box_h)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = C["dark"]
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top    = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    tf.margin_left   = Inches(0.2)
    tf.margin_right  = Inches(0.2)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = hero_label
    run.font.name = BRAND["font"]
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = C["white"]

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = hero_curr
    run2.font.name = BRAND["font"]
    run2.font.size = Pt(52)
    run2.font.bold = True
    run2.font.color.rgb = C["gold"]

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = f"vs {hero_prev}  ({hero_pct})"
    run3.font.name = BRAND["font"]
    run3.font.size = Pt(11)
    run3.font.color.rgb = C["white"]

    _add_chart_image(slide, chart_path,
                     left=Inches(4.7), top=box_top, width=Inches(5.0))
    return slide


def slide_scorecard(prs, title, kpis):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
    _add_title_textbox(slide, title)

    n = len(kpis)
    box_w  = Inches(2.5)
    box_h  = Inches(1.1)
    gap    = Inches(0.2)
    total_w = n * box_w + (n - 1) * gap
    start_x = (prs.slide_width - total_w) // 2
    start_y = (prs.slide_height - box_h) // 2
    box_colours = [C["gold"], C["orange"], C["teal"]]

    for i, (label, data) in enumerate(kpis):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(5, x, start_y, box_w, box_h)
        shape.adjustments[0] = 0.08
        shape.fill.solid()
        shape.fill.fore_color.rgb = box_colours[i % len(box_colours)]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_top    = Inches(0.1)
        tf.margin_bottom = Inches(0.08)
        tf.margin_left   = Inches(0.1)
        tf.margin_right  = Inches(0.1)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = BRAND["font"]
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = C["dark"]

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = data.get('curr', '—')
        run2.font.name = BRAND["font"]
        run2.font.size = Pt(17)
        run2.font.bold = True
        run2.font.color.rgb = C["dark"]

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = f"vs {data.get('prev', '—')}  ({data.get('pct', '—')})"
        run3.font.name = BRAND["font"]
        run3.font.size = Pt(9)
        run3.font.color.rgb = C["dark"]

    return slide


def slide_table(prs, title, headers, rows, status_col=None, totals_row=None):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
    _add_title_textbox(slide, title)
    table_top = Inches(1.2)
    table_h = prs.slide_height - table_top - Inches(0.2)
    _add_table_shape(slide, headers, rows,
                     left=Inches(0.5), top=table_top,
                     width=prs.slide_width - Inches(1.0), height=table_h,
                     status_col=status_col, totals_row=totals_row)
    return slide


def slide_table_commentary(prs, title, headers, rows, bullets, status_col=None, totals_row=None):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
    _add_title_textbox(slide, title)

    table_top = Inches(1.2)
    content_h = prs.slide_height - table_top - Inches(0.2)

    tb = slide.shapes.add_textbox(Inches(0.5), table_top, Inches(4.0), content_h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        text = bullet['point'] if isinstance(bullet, dict) else bullet
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        _shrink_bullet(p)
        for run in p.runs:
            run.font.size = Pt(11)
            run.font.name = BRAND["font"]
            run.font.color.rgb = C["dark"]

    _add_table_shape(slide, headers, rows,
                     left=Inches(4.7), top=table_top,
                     width=prs.slide_width - Inches(5.0), height=content_h,
                     status_col=status_col, totals_row=totals_row)
    return slide


def slide_planning_gantt(prs, title, tasks, plan_start_str=None, plan_end_str=None):
    if not tasks:
        return None

    parsed = []
    for t in tasks:
        start_str = _fmt_date(t.get('start_date', ''))
        end_str   = _fmt_date(t.get('end_date', ''))
        try:
            start = datetime.strptime(start_str, '%d/%m/%Y')
            end   = datetime.strptime(end_str,   '%d/%m/%Y')
        except ValueError:
            continue
        platform = t.get('platform', '')
        name     = t.get('name', '')
        label    = f"{platform}: {name}" if platform and name else platform or name
        parsed.append({'label': label, 'start': start, 'end': end,
                        'status': t.get('status', 'Scheduled')})

    if not parsed:
        return None

    # Anchor the time axis to the plan window when available, else fall back to task span
    def _parse_plan_date(s):
        if not s:
            return None
        for fmt in ('%d/%m/%y', '%d/%m/%Y', '%Y-%m-%d'):
            try:
                return datetime.strptime(s, fmt)
            except ValueError:
                continue
        return None

    plan_start_dt = _parse_plan_date(plan_start_str)
    plan_end_dt   = _parse_plan_date(plan_end_str)

    period_start = (plan_start_dt if plan_start_dt else min(t['start'] for t in parsed)).replace(day=1)
    axis_end     = plan_end_dt if plan_end_dt else max(t['end'] for t in parsed)
    last_month_start = axis_end.replace(day=1)

    months = []
    m = period_start
    while m <= last_month_start:
        months.append(m)
        m = (m.replace(month=m.month + 1) if m.month < 12
             else m.replace(year=m.year + 1, month=1))

    last_m     = months[-1]
    period_end = last_m.replace(day=monthrange(last_m.year, last_m.month)[1])
    total_days = (period_end - period_start).days + 1

    # Layout constants (EMU)
    SL = prs.slide_width
    SH = prs.slide_height
    LEFT         = Inches(0.4)
    RIGHT        = Inches(0.4)
    TITLE_BOTTOM = Inches(1.1)
    BOTTOM       = Inches(0.3)
    LABEL_W      = Inches(2.5)
    CHART_LEFT   = LEFT + LABEL_W + Inches(0.1)
    CHART_W      = SL - CHART_LEFT - RIGHT
    HEADER_H     = Inches(0.35)
    LEGEND_H     = Inches(0.35)
    LEGEND_PAD   = Inches(0.12)
    available    = SH - TITLE_BOTTOM - BOTTOM - HEADER_H - LEGEND_H - LEGEND_PAD
    row_h        = max(Inches(0.3), min(Inches(0.5), available // len(parsed)))
    ROWS_TOP     = TITLE_BOTTOM + HEADER_H

    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
    _add_title_textbox(slide, title)

    def _rect(left, top, width, height, fill=None):
        shp = slide.shapes.add_shape(
            1, int(left), int(top), int(max(width, 1)), int(max(height, 1))
        )
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        else:
            shp.fill.background()
        sp_pr = shp._element.find(qn('p:spPr'))
        if sp_pr is not None:
            ln = sp_pr.find(qn('a:ln'))
            if ln is None:
                ln = etree.SubElement(sp_pr, qn('a:ln'))
            else:
                for child in list(ln):
                    ln.remove(child)
            etree.SubElement(ln, qn('a:noFill'))
        return shp

    def _label(shp, text, size=Pt(9), bold=False, color=C["dark"], align=PP_ALIGN.LEFT):
        tf = shp.text_frame
        tf.word_wrap        = True
        tf.vertical_anchor  = MSO_ANCHOR.MIDDLE
        tf.margin_left      = Inches(0.06)
        tf.margin_right     = Inches(0.04)
        tf.margin_top       = Inches(0.02)
        tf.margin_bottom    = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text           = text
        run.font.name      = BRAND["font"]
        run.font.size      = size
        run.font.bold      = bold
        run.font.color.rgb = color

    # Month header cells
    for i, month in enumerate(months):
        offset    = (month - period_start).days
        span      = ((months[i + 1] - month).days if i < len(months) - 1
                     else (period_end - month).days + 1)
        col_left  = CHART_LEFT + int(CHART_W * offset / total_days)
        col_w     = int(CHART_W * span / total_days)
        shp = _rect(col_left, TITLE_BOTTOM, col_w, HEADER_H, fill=C["dark"])
        _label(shp, month.strftime('%B'), size=Pt(10), bold=True,
               color=C["white"], align=PP_ALIGN.CENTER)

    # Task rows
    for i, task in enumerate(parsed):
        row_top = ROWS_TOP + i * row_h
        _rect(CHART_LEFT, row_top, CHART_W, row_h,
              fill=C["light"] if i % 2 == 0 else C["white"])

        label_shp = _rect(LEFT, row_top, LABEL_W, row_h, fill=C["dark"])
        _label(label_shp, task['label'], size=Pt(8), color=C["white"])

        s = max(task['start'], period_start)
        e = min(task['end'],   period_end)
        if s > e:
            continue
        pad   = row_h // 6
        s_off = (s - period_start).days
        e_off = (e - period_start).days + 1
        bar_l = CHART_LEFT + int(CHART_W * s_off / total_days)
        bar_w = max(int(CHART_W * (e_off - s_off) / total_days), Inches(0.05))
        _rect(bar_l, row_top + pad, bar_w, row_h - 2 * pad,
              fill=STATUS_COLOURS.get(task['status'], C["grey"]))

    # Legend (centered below rows)
    legend_top = ROWS_TOP + len(parsed) * row_h + LEGEND_PAD
    swatch     = Inches(0.14)
    text_w_l   = Inches(0.9)
    spacing    = Inches(0.2)
    item_w     = swatch + Inches(0.07) + text_w_l
    items      = [('Complete',    C["teal"]),  ('In Progress', C["gold"]),
                  ('Scheduled',   C["grey"]),  ('Blocked',     C["orange"])]
    total_lw   = len(items) * item_w + (len(items) - 1) * spacing
    lx0        = (SL - total_lw) // 2

    for j, (lbl, col) in enumerate(items):
        lx = lx0 + j * (item_w + spacing)
        ly = legend_top + (LEGEND_H - swatch) // 2
        _rect(lx, ly, swatch, swatch, fill=col)
        tb  = slide.shapes.add_textbox(lx + swatch + Inches(0.07), legend_top, text_w_l, LEGEND_H)
        tf  = tb.text_frame
        run = tf.paragraphs[0].add_run()
        run.text           = lbl
        run.font.name      = BRAND["font"]
        run.font.size      = Pt(8)
        run.font.color.rgb = C["dark"]

    return slide


KANBAN_COLUMNS = ['Blocked', 'Scheduled', 'Complete']
KANBAN_HEADER_COLOUR = {
    'Blocked':   C['gold'],
    'Scheduled': C['orange'],
    'Complete':  C['teal'],
}
KANBAN_HEADER_TEXT = {
    'Blocked':   C['dark'],
    'Scheduled': C['white'],
    'Complete':  C['white'],
}


def slide_action_kanban(prs, title, tasks):
    if not tasks:
        return None

    grouped = {s: [] for s in KANBAN_COLUMNS}
    for task in tasks:
        status = task.get('status', 'Scheduled')
        # In Progress folds into Scheduled
        bucket = status if status in grouped else 'Scheduled'
        grouped[bucket].append(task)

    for status in KANBAN_COLUMNS:
        grouped[status].sort(key=lambda t: _parse_date_for_sort(t.get('end_date', '')))

    n_cols     = len(KANBAN_COLUMNS)
    margin_x   = Inches(0.4)
    margin_top = Inches(1.1)
    margin_bot = Inches(0.3)
    col_gap    = Inches(0.2)
    header_h   = Inches(0.4)
    card_gap   = Inches(0.15)

    col_w            = (prs.slide_width - 2 * margin_x - (n_cols - 1) * col_gap) // n_cols
    available_card_h = prs.slide_height - margin_top - margin_bot - header_h - card_gap
    max_tasks        = max((len(grouped[s]) for s in KANBAN_COLUMNS), default=1) or 1
    card_h           = min(Inches(1.1), max(Inches(0.7),
                           (available_card_h - (max_tasks - 1) * card_gap) // max_tasks))

    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
    _add_title_textbox(slide, title)

    for col_i, status in enumerate(KANBAN_COLUMNS):
        col_x = margin_x + col_i * (col_w + col_gap)

        # Rounded column header
        hdr = slide.shapes.add_shape(5, int(col_x), int(margin_top), int(col_w), int(header_h))
        hdr.adjustments[0] = 0.08
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = KANBAN_HEADER_COLOUR[status]
        hdr.line.fill.background()
        tf = hdr.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = status
        run.font.name = BRAND["font"]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = KANBAN_HEADER_TEXT[status]

        for card_i, task in enumerate(grouped[status]):
            card_y = margin_top + header_h + card_gap + card_i * (card_h + card_gap)

            # Rounded dark card
            card = slide.shapes.add_shape(5, int(col_x), int(card_y), int(col_w), int(card_h))
            card.adjustments[0] = 0.08
            card.fill.solid()
            card.fill.fore_color.rgb = C["dark"]
            card.line.fill.background()

            pad = Inches(0.15)
            tb = slide.shapes.add_textbox(
                int(col_x + pad), int(card_y + pad),
                int(col_w - 2 * pad), int(card_h - 2 * pad)
            )
            tf = tb.text_frame
            tf.word_wrap = True

            platform = task.get('platform', '')
            p0 = tf.paragraphs[0]
            if platform:
                run0 = p0.add_run()
                run0.text = platform
                run0.font.name = BRAND["font"]
                run0.font.size = Pt(7)
                run0.font.bold = True
                run0.font.color.rgb = C["gold"]

            p1 = tf.add_paragraph()
            run1 = p1.add_run()
            run1.text = task.get('name', '')
            run1.font.name = BRAND["font"]
            run1.font.size = Pt(9)
            run1.font.bold = True
            run1.font.color.rgb = C["white"]

            end_date = _fmt_date(task.get('end_date', ''))
            if end_date:
                date_label = 'Completed' if status == 'Complete' else 'Due'
                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = f'{date_label}: {end_date}'
                run2.font.name = BRAND["font"]
                run2.font.size = Pt(7)
                run2.font.color.rgb = C["grey"]

    return slide


# ── TABLE DATA EXTRACTION ────────────────────────────────────────────────────

_ADDITIVE_METRICS = {
    'Cost', 'Transaction Revenue', 'Conversions', 'Impressions',
    'Clicks', 'Transactions', 'Views', 'Hooks', 'Holds',
}

# (numerator, denominator, multiplier, format_type)
_DERIVED_METRIC_FORMULAS = {
    'ROAS':            ('Transaction Revenue', 'Cost',               100, 'pct'),
    'CPA':             ('Cost',                'Conversions',           1, 'gbp'),
    'CTR':             ('Clicks',              'Impressions',         100, 'pct'),
    'CPC':             ('Cost',                'Clicks',                1, 'gbp'),
    'Conversion Rate': ('Conversions',         'Clicks',              100, 'pct'),
    'AOV':             ('Transaction Revenue', 'Transactions',          1, 'gbp'),
    'View Rate':       ('Views',               'Impressions',         100, 'pct'),
    'Hook Rate':       ('Hooks',               'Impressions',         100, 'pct'),
    'Hold Rate':       ('Holds',               'Impressions',         100, 'pct'),
    'Cost Per Hook':   ('Cost',                'Hooks',                 1, 'gbp'),
}

_FMT_INT_METRICS = {
    'Conversions', 'Impressions', 'Clicks', 'Transactions', 'Views', 'Hooks', 'Holds',
}


def _parse_num(s):
    try:
        return float(str(s).replace('£', '').replace('%', '').replace(',', '').replace('x', '').strip())
    except (ValueError, AttributeError):
        return 0.0


def _fmt_metric(metric, value):
    if metric in _FMT_INT_METRICS:
        return fmt_int(value)
    elif metric in ('Cost', 'Transaction Revenue') or metric in _DERIVED_METRIC_FORMULAS and _DERIVED_METRIC_FORMULAS[metric][3] == 'gbp':
        return fmt_gbp(value)
    else:
        return fmt_pct(value)


def _apply_row_filters(items, row_filters, dimension_col):
    def _passes(dim_val, metric_dict):
        for f in row_filters:
            col = f.get('column', '')
            op  = f.get('op', '=')
            val = f.get('value')
            if col == dimension_col:
                s = str(dim_val)
                sv = str(val)
                if   op == 'contains'     and sv not in s:  return False
                elif op == 'not_contains' and sv in s:       return False
                elif op == '='            and s != sv:       return False
                elif op == '!='           and s == sv:       return False
            else:
                raw = metric_dict.get(col, {})
                curr_str = raw.get('curr', '0') if isinstance(raw, dict) else '0'
                try:
                    n  = _parse_num(curr_str)
                    fv = float(val)
                    if   op == '>'  and not (n >  fv): return False
                    elif op == '<'  and not (n <  fv): return False
                    elif op == '>=' and not (n >= fv): return False
                    elif op == '<=' and not (n <= fv): return False
                    elif op == '='  and not (n == fv): return False
                    elif op == '!=' and not (n != fv): return False
                except (ValueError, TypeError):
                    return False
        return True
    return [(dv, md) for dv, md in items if _passes(dv, md)]


def _compute_metric_raw(metric, sums):
    if metric in _ADDITIVE_METRICS:
        return sums.get(metric, 0.0)
    if metric in _DERIVED_METRIC_FORMULAS:
        num_key, denom_key, mult, _ = _DERIVED_METRIC_FORMULAS[metric]
        denom = sums.get(denom_key, 0.0)
        if denom == 0:
            return None
        return (sums.get(num_key, 0.0) / denom) * mult
    return None


def _compute_totals_row(items, metrics, comparison):
    needed = set()
    for m in metrics:
        if m in _ADDITIVE_METRICS:
            needed.add(m)
        elif m in _DERIVED_METRIC_FORMULAS:
            num_key, denom_key, _, _ = _DERIVED_METRIC_FORMULAS[m]
            needed.add(num_key)
            needed.add(denom_key)

    curr_sums = {k: 0.0 for k in needed}
    prev_sums = {k: 0.0 for k in needed} if comparison else {}

    for _, metric_dict in items:
        for comp in needed:
            vals = metric_dict.get(comp, {})
            if isinstance(vals, dict):
                curr_sums[comp] += _parse_num(vals.get('curr', '0'))
                if comparison:
                    prev_sums[comp] += _parse_num(vals.get('prev', '0'))

    row = ['Totals']
    for m in metrics:
        curr_raw = _compute_metric_raw(m, curr_sums)
        row.append(_fmt_metric(m, curr_raw) if curr_raw is not None else '—')
        if comparison:
            prev_raw = _compute_metric_raw(m, prev_sums)
            row.append(_fmt_metric(m, prev_raw) if prev_raw is not None else '—')
            if curr_raw is not None and prev_raw and prev_raw != 0:
                pct = (curr_raw - prev_raw) / abs(prev_raw) * 100
                row.append(f"+{pct:.1f}%" if pct >= 0 else f"{pct:.1f}%")
            else:
                row.append('—')
    return row


def render_table_data(graph, client, max_rows=12, comparison=True):
    """Convert dimension_data into (headers, rows, totals_row) for table slides.

    comparison=True  (graph_type 'table_comparison'): curr + prev + % change per metric.
    comparison=False (graph_type 'table'):             current period values only.
    Returns (headers, rows, totals_row) — totals_row is None when show_totals is falsy.
    """
    data_source = graph.get('data_source')
    if not data_source:
        return [], [], None

    dim_entry     = client.get('dimension_data', {}).get(data_source, {})
    comp_key      = 'yoy' if graph.get('comparison') == 'yoy' else 'mom'
    comp_data     = dim_entry.get(comp_key, {})
    metrics       = graph.get('metrics', [])
    dimension_col = data_source.split('::')[0]
    row_filters   = graph.get('row_filters', [])
    sort_by       = graph.get('sort_by')
    sort_dir      = graph.get('sort_dir', 'desc')
    show_totals   = graph.get('show_totals', False)

    if not comp_data or not metrics:
        return [], [], None

    items = [(dv, md) for dv, md in comp_data.items() if isinstance(md, dict)]

    if row_filters:
        items = _apply_row_filters(items, row_filters, dimension_col)

    totals_row = _compute_totals_row(items, metrics, comparison) if show_totals else None

    rows = []
    for dim_val, metric_dict in items:
        row = [str(dim_val)]
        for m in metrics:
            vals = metric_dict.get(m, {})
            if isinstance(vals, dict):
                row.append(vals.get('curr', '—'))
                if comparison:
                    row.append(vals.get('prev', '—'))
                    row.append(vals.get('pct',  '—'))
            else:
                row.append('—')
                if comparison:
                    row.extend(['—', '—'])
        rows.append(row)

    # Sort: find column index for sort_by, default to first metric col
    if sort_by and sort_by in metrics:
        col_idx = 1 + metrics.index(sort_by) * (3 if comparison else 1)
    else:
        col_idx = 1

    def _sort_key(r):
        return _parse_num(r[col_idx]) if col_idx < len(r) else 0.0

    rows.sort(key=_sort_key, reverse=(sort_dir != 'asc'))
    rows = rows[:max_rows]

    headers = [dimension_col]
    for m in metrics:
        headers.append(m)
        if comparison:
            headers += [f"{m} (prev)", f"{m} (%)"]

    return headers, rows, totals_row


# ── SLIDE TEMPLATE REGISTRY ───────────────────────────────────────────────────

SLIDE_TEMPLATE_REGISTRY = {
    'chart_commentary',
    'full_chart',
    'big_number',
    'scorecard_vertical',
    'scorecard_horizontal',
    'table_commentary',
    'table',
}


# ── TEMPLATE DISPATCH HELPERS ────────────────────────────────────────────────

def _resolve_date_label_for_graph(client, graph):
    _rd   = client.get('dimension_data', {}).get(graph.get('data_source', ''), {}).get('resolved_dates', {})
    _comp = graph.get('comparison')
    if _comp == 'mom':
        comp_start, comp_end = _rd.get('prev_start'), _rd.get('prev_end')
    elif _comp == 'yoy':
        comp_start, comp_end = _rd.get('yoy_start'), _rd.get('yoy_end')
    else:
        comp_start, comp_end = None, None
    return _fmt_date_label(
        _rd.get('current_start', graph['date_range']['start']),
        _rd.get('current_end',   graph['date_range']['end']),
        comp_start, comp_end,
    )


def _extract_hero_metric(client, graph):
    """Return (label, curr, prev, pct) for the first metric from mom data."""
    data_source = graph.get('data_source', '')
    dim_entry   = client.get('dimension_data', {}).get(data_source, {})
    mom_data    = dim_entry.get('mom', {})
    metric      = (graph.get('metrics') or [''])[0]

    totals = {}
    for dim_val, metric_dict in mom_data.items():
        if not isinstance(metric_dict, dict):
            continue
        vals = metric_dict.get(metric, {})
        if isinstance(vals, dict):
            totals = vals
            break

    return (
        metric,
        totals.get('curr', '—'),
        totals.get('prev', '—'),
        totals.get('pct',  '—'),
    )


def _render_overview_slide(prs, client, template, title, summary, bullets, kpis, date_label):
    if template == 'scorecard_horizontal':
        slide_scorecard_horizontal(prs, title, summary, bullets, kpis, date_label)
    elif template == 'chart_commentary':
        slide_commentary(prs, title, summary, bullets)
    else:
        slide_scorecard_vertical(prs, title, summary, bullets, kpis, date_label)


def _render_trend_slide(prs, client, trend):
    template   = trend.get('template', 'chart_commentary')
    graph      = trend.get('graph', {})
    title      = trend['title']
    summary    = trend['summary']
    bullets    = trend['bullets']
    date_label = _resolve_date_label_for_graph(client, graph) if graph.get('data_source') else None

    if template in ('table', 'table_commentary'):
        comparison = graph.get('graph_type') == 'table_comparison'
        headers, rows, totals_row = render_table_data(graph, client, comparison=comparison)
        if not headers:
            slide_commentary(prs, title, summary, bullets)
            return
        if template == 'table_commentary':
            slide_table_commentary(prs, title, headers, rows, bullets, totals_row=totals_row)
        else:
            slide_table(prs, title, headers, rows, totals_row=totals_row)
        return

    chart_path = render_graph(client, graph) if graph.get('data_source') else None

    if template == 'full_chart':
        if chart_path:
            slide_full_chart(prs, title, summary, chart_path, date_label)
        else:
            slide_commentary(prs, title, summary, bullets)

    elif template == 'big_number':
        if chart_path:
            hero_label, hero_curr, hero_prev, hero_pct = _extract_hero_metric(client, graph)
            slide_big_number(prs, title, summary, chart_path,
                             hero_label, hero_curr, hero_prev, hero_pct, date_label)
        else:
            slide_commentary(prs, title, summary, bullets)

    elif template in ('scorecard_vertical', 'scorecard_horizontal'):
        # Scorecard as trend: build KPI boxes from top dimension values in mom data
        data_source = graph.get('data_source', '')
        dim_entry   = client.get('dimension_data', {}).get(data_source, {})
        mom_data    = dim_entry.get('mom', {})
        metric      = (graph.get('metrics') or [''])[0]
        kpi_items = []
        for dim_val, metric_dict in list(mom_data.items())[:4]:
            if not isinstance(metric_dict, dict):
                continue
            vals = metric_dict.get(metric, {})
            if isinstance(vals, dict):
                kpi_items.append((str(dim_val), vals))
        if kpi_items:
            if template == 'scorecard_horizontal':
                slide_scorecard_horizontal(prs, title, summary, bullets, kpi_items, date_label)
            else:
                slide_scorecard_vertical(prs, title, summary, bullets, kpi_items, date_label)
        else:
            slide_commentary(prs, title, summary, bullets)

    else:
        # Default: chart_commentary
        if chart_path:
            slide_chart_commentary(prs, title, summary, bullets, chart_path, date_label)
        else:
            slide_commentary(prs, title, summary, bullets)


# ── PIPELINE ORCHESTRATOR ─────────────────────────────────────────────────────

def generate_ppt(client_name, output_path=None, slide_content=None):
    data_path = os.path.join(PROJECT_ROOT, "storage", f"{client_name}_monthly_data.json")
    if not os.path.exists(data_path):
        raise FileNotFoundError(
            f"No monthly data file found for '{client_name}'. "
            f"Run: python -m monthly_reports.main --client \"{client_name}\""
        )
    with open(data_path, "r", encoding="utf-8") as f:
        client = json.load(f)

    if slide_content is None:
        slide_content = generate_monthly_slide_content(client)
        mtd_content = generate_mtd_slide_content(client)
        slide_content = {**slide_content, **mtd_content}
    client['slide_content'] = slide_content
    content_path = os.path.join(PROJECT_ROOT, "storage", f"{client_name}_monthly_content.json")
    with open(content_path, "w", encoding="utf-8") as f:
        json.dump(client['slide_content'], f, ensure_ascii=False, indent=2)

    prev_month = datetime.strptime(client['start_date_string'], "%d/%m/%Y").strftime("%B")

    if output_path is None:
        timestamp = datetime.now().strftime("%Y-%m-%d_%H%M%S")
        output_path = os.path.join(PROJECT_ROOT, "slides", f"{client_name}_monthly_{timestamp}.pptx")
    template_path = os.path.join(PROJECT_ROOT, "slides", "template.pptx")
    shutil.copy(template_path, output_path)
    prs = Presentation(output_path)

    slide_rid = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(slide_rid)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    sc = client['slide_content']

    slide_cover(prs, f'{client["name"]} Monthly Deck')

    slide_section_separator(prs, 'Paid Media', variant='navy')

    # ── Performance Overview ──────────────────────────────────────────────────
    overview_template = sc['overview'].get('template', 'scorecard_vertical')
    overview_kpi_count = sc['overview'].get('kpi_count', None)
    kpis = _build_kpis_for(client, 'paid_data', overview_kpi_count)
    overview_date_label = _fmt_date_label(
        client['start_date_string'],
        client['end_date_string'],
        _iso_to_dmy(client.get('compare_start_mom', '')),
        _iso_to_dmy(client.get('compare_end_mom', '')),
    )
    slide_section_separator(prs, f'{prev_month} Performance Overview', variant='gold')
    _render_overview_slide(
        prs, client, overview_template,
        title=f'{prev_month} Performance',
        summary=sc['overview']['summary'],
        bullets=sc['overview']['bullets'],
        kpis=kpis,
        date_label=overview_date_label,
    )

    # ── MTD Overview ─────────────────────────────────────────────────────────
    mtd_start_str = client.get('mtd_start_date_string')
    if mtd_start_str and sc.get('mtd_overview'):
        mtd_month = datetime.strptime(mtd_start_str, "%d/%m/%Y").strftime("%B")
        mtd_template = sc['mtd_overview'].get('template', 'scorecard_vertical')
        mtd_kpi_count = sc['mtd_overview'].get('kpi_count', None)
        mtd_kpis = _build_kpis_for(client, 'paid_data_mtd', mtd_kpi_count)
        mtd_date_label = _fmt_date_label(
            client['mtd_start_date_string'],
            client['mtd_end_date_string'],
            _iso_to_dmy(client.get('compare_start_mtd', '')),
            _iso_to_dmy(client.get('compare_end_mtd', '')),
        )
        slide_section_separator(prs, f'{mtd_month} Performance Overview', variant='gold')
        _render_overview_slide(
            prs, client, mtd_template,
            title=f'{mtd_month} Performance',
            summary=sc['mtd_overview']['summary'],
            bullets=sc['mtd_overview']['bullets'],
            kpis=mtd_kpis,
            date_label=mtd_date_label,
        )

    # ── Trend Slides ─────────────────────────────────────────────────────────
    slide_section_separator(prs, 'Top Level Trends', variant='gold')
    for trend in sc['trends']:
        _render_trend_slide(prs, client, trend)

    slide_section_separator(prs, 'Plan Overview', variant='gold')
    plan_json = client.get('plan_json')
    current_tasks = _extract_current_tasks(plan_json) if plan_json else []
    if current_tasks:
        slide_action_kanban(prs, 'Plan Overview', current_tasks)
    all_tasks, plan_start, plan_end = _extract_all_plan_tasks(plan_json) if plan_json else ([], None, None)
    if all_tasks:
        slide_planning_gantt(prs, '90 Day Plan', all_tasks, plan_start, plan_end)

    prs.save(output_path)
    print(f"Saved {output_path} with {len(prs.slides)} slide(s)")

    excel_path = export_slide_data(client, sc, output_path)
    if excel_path:
        print(f"Saved data export to {excel_path}")

    return output_path, excel_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--client", required=True, help="Client name as it appears in config.json")
    parser.add_argument("--output", default=None, help="Output path for the generated PPTX")
    args = parser.parse_args()
    generate_ppt(args.client, args.output)
