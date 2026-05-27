import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pytest
from pptx import Presentation
from pptx.util import Inches

from monthly_reports.generate_ppt import (
    C,
    STATUS_COLOURS,
    SLD_LAYOUT_BLANK,
    SLD_LAYOUT_TITLE_AND_BODY,
    _extract_current_tasks,
    _add_table_shape,
    _fmt_date,
    slide_planning_gantt,
    slide_scorecard_commentary,
)


@pytest.fixture
def prs():
    p = Presentation('slides/template.pptx')
    slide_rid = p.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    p.part.drop_rel(slide_rid)
    p.slides._sldIdLst.remove(p.slides._sldIdLst[0])
    return p


# ── _extract_current_tasks ────────────────────────────────────────────────────

class TestExtractCurrentTasks:
    def test_returns_tasks_for_current_quarter(self):
        plan = {
            "Q1 2026": {"plan_status": "old",     "tasks": [{"name": "Old task"}]},
            "Q2 2026": {"plan_status": "current",  "tasks": [{"name": "Current task"}]},
        }
        assert _extract_current_tasks(plan) == [{"name": "Current task"}]

    def test_returns_empty_when_no_current_quarter(self):
        plan = {"Q1 2026": {"plan_status": "old", "tasks": [{"name": "Old task"}]}}
        assert _extract_current_tasks(plan) == []

    def test_returns_empty_for_none(self):
        assert _extract_current_tasks(None) == []

    def test_returns_empty_for_empty_dict(self):
        assert _extract_current_tasks({}) == []

    def test_returns_empty_tasks_list_when_current_quarter_has_no_tasks(self):
        plan = {"Q2 2026": {"plan_status": "current", "tasks": []}}
        assert _extract_current_tasks(plan) == []

    def test_handles_flat_list_shape(self):
        tasks = [{"name": "Task A"}, {"name": "Task B"}]
        assert _extract_current_tasks(tasks) == tasks

    def test_handles_tasks_keyed_shape(self):
        plan = {"tasks": [{"name": "Task A"}]}
        assert _extract_current_tasks(plan) == [{"name": "Task A"}]

    def test_picks_current_among_multiple_quarters(self):
        plan = {
            "Q4 2025": {"plan_status": "old",     "tasks": [{"name": "Old"}]},
            "Q1 2026": {"plan_status": "old",     "tasks": [{"name": "Also old"}]},
            "Q2 2026": {"plan_status": "current",  "tasks": [{"name": "Current"}]},
        }
        result = _extract_current_tasks(plan)
        assert len(result) == 1
        assert result[0]["name"] == "Current"


# ── _fmt_date ─────────────────────────────────────────────────────────────────

class TestFmtDate:
    def test_two_digit_year(self):
        assert _fmt_date("06/04/26") == "06/04/2026"

    def test_four_digit_year(self):
        assert _fmt_date("06/04/2026") == "06/04/2026"

    def test_iso_format(self):
        assert _fmt_date("2026-04-06") == "06/04/2026"

    def test_empty_string(self):
        assert _fmt_date("") == ""

    def test_none_returns_empty(self):
        assert _fmt_date(None) == ""


# ── _add_table_shape ──────────────────────────────────────────────────────────

class TestAddTableShape:
    def test_creates_table_with_correct_row_count(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        headers = ['Name', 'Status']
        rows    = [['Task A', 'Complete'], ['Task B', 'Scheduled']]
        _add_table_shape(slide, headers, rows, Inches(0.5), Inches(1), Inches(9), Inches(4))

        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1
        tbl = tables[0].table
        assert len(tbl.rows) == 3  # 1 header + 2 data rows

    def test_header_cell_text(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name', 'Status'], [['Task A', 'Complete']],
                         Inches(0.5), Inches(1), Inches(9), Inches(2))
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.rows[0].cells[0].text == 'Name'
        assert tbl.rows[0].cells[1].text == 'Status'

    def test_data_cell_text(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name'], [['Task A']], Inches(0.5), Inches(1), Inches(9), Inches(2))
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.rows[1].cells[0].text == 'Task A'

    def test_header_uses_dark_background(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name'], [['Task A']], Inches(0.5), Inches(1), Inches(9), Inches(2))
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.cell(0, 0).fill.fore_color.rgb == C["dark"]

    def test_status_col_applies_status_colour_for_complete(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name', 'Status'], [['Task A', 'Complete']],
                         Inches(0.5), Inches(1), Inches(9), Inches(2), status_col=1)
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.cell(1, 1).fill.fore_color.rgb == STATUS_COLOURS['Complete']

    def test_status_col_applies_status_colour_for_blocked(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name', 'Status'], [['Task B', 'Blocked']],
                         Inches(0.5), Inches(1), Inches(9), Inches(2), status_col=1)
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.cell(1, 1).fill.fore_color.rgb == STATUS_COLOURS['Blocked']

    def test_alternating_row_shading(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        rows = [['A', 'Complete'], ['B', 'Complete'], ['C', 'Complete']]
        _add_table_shape(slide, ['Name', 'Status'], rows,
                         Inches(0.5), Inches(1), Inches(9), Inches(4))
        tbl = [s for s in slide.shapes if s.has_table][0].table
        # Row 1 (i=0) → light; row 2 (i=1) → white; row 3 (i=2) → light
        assert tbl.cell(1, 0).fill.fore_color.rgb == C["light"]
        assert tbl.cell(2, 0).fill.fore_color.rgb == C["white"]
        assert tbl.cell(3, 0).fill.fore_color.rgb == C["light"]


# ── slide_planning_gantt ──────────────────────────────────────────────────────

class TestSlidePlanningGantt:
    TASKS = [
        {'name': 'Task A', 'platform': 'Google Ads',
         'start_date': '06/04/26', 'end_date': '10/04/26', 'status': 'Complete'},
        {'name': 'Task B', 'platform': 'Meta Ads',
         'start_date': '12/04/26', 'end_date': '25/04/26', 'status': 'Scheduled'},
    ]

    def test_returns_none_for_empty_tasks(self, prs):
        assert slide_planning_gantt(prs, 'Q2 Plan', []) is None

    def test_returns_slide_for_valid_tasks(self, prs):
        slide = slide_planning_gantt(prs, 'Q2 Plan', self.TASKS)
        assert slide is not None

    def test_slide_contains_image_not_table(self, prs):
        slide = slide_planning_gantt(prs, 'Q2 Plan', self.TASKS)
        pictures = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE = 13
        tables   = [s for s in slide.shapes if s.has_table]
        assert len(pictures) >= 1
        assert len(tables) == 0

    def test_gantt_image_file_is_created(self, prs):
        import os
        slide_planning_gantt(prs, 'Q2 Plan Test', self.TASKS)
        assert os.path.exists('charts/Q2_Plan_Test_gantt.png')


# ── slide_scorecard_commentary ────────────────────────────────────────────────

class TestSlideScorecardCommentary:
    def _make_kpis(self):
        return [
            ('Cost',    {'curr': '£1,000', 'prev': '£900',  'pct': '+11%'}),
            ('Revenue', {'curr': '£5,000', 'prev': '£4,000', 'pct': '+25%'}),
            ('ROAS',    {'curr': '5.0',    'prev': '4.5',    'pct': '+11%'}),
        ]

    def test_title_placeholder_is_set(self, prs):
        slide = slide_scorecard_commentary(
            prs, 'Top Level View', 'Good month', [{'point': 'Revenue up'}], self._make_kpis()
        )
        assert slide.placeholders[0].text == 'Top Level View'

    def test_kpi_boxes_are_added(self, prs):
        kpis = self._make_kpis()
        slide = slide_scorecard_commentary(
            prs, 'Top Level View', 'Good month', [{'point': 'Revenue up'}], kpis
        )
        # Shapes include layout placeholders + 3 KPI boxes; at least 3 shapes total
        assert len(slide.shapes) >= len(kpis)

    def test_returns_slide_object(self, prs):
        slide = slide_scorecard_commentary(
            prs, 'Top Level View', 'Good month', [{'point': 'Revenue up'}], self._make_kpis()
        )
        assert slide is not None
